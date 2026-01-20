import asyncio
import tempfile
from pathlib import Path
import logging
from typing import List, Optional
import os

logger = logging.getLogger(__name__)

class PDFConverter:
    """Handle PDF conversion to/from various formats"""
    
    def __init__(self):
        self.supported_formats = {
            "pdf": [".pdf"],
            "word": [".docx", ".doc"],
            "excel": [".xlsx", ".xls"],
            "powerpoint": [".pptx", ".ppt"],
            "image": [".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".gif"],
            "text": [".txt", ".html", ".rtf"]
        }
    
    async def convert(self, input_path: Path, output_format: str, 
                     quality: str = "high", pages: Optional[List[int]] = None) -> Path:
        """Convert file to specified format"""
        
        # Check if input is PDF
        if input_path.suffix.lower() == ".pdf":
            return await self._convert_from_pdf(input_path, output_format, quality, pages)
        else:
            return await self._convert_to_pdf(input_path, output_format, quality)
    
    async def _convert_from_pdf(self, pdf_path: Path, output_format: str,
                               quality: str, pages: Optional[List[int]]) -> Path:
        """Convert PDF to other formats"""
        
        output_dir = Path("processed")
        output_name = f"{pdf_path.stem}_converted.{output_format}"
        output_path = output_dir / output_name
        
        try:
            if output_format in ["docx", "doc"]:
                return await self._pdf_to_word(pdf_path, output_path, pages)
            
            elif output_format in ["xlsx", "xls"]:
                return await self._pdf_to_excel(pdf_path, output_path, pages)
            
            elif output_format in ["pptx", "ppt"]:
                return await self._pdf_to_powerpoint(pdf_path, output_path, pages)
            
            elif output_format in ["jpg", "jpeg", "png", "tiff"]:
                return await self._pdf_to_image(pdf_path, output_path, output_format, quality, pages)
            
            elif output_format in ["txt", "html"]:
                return await self._pdf_to_text(pdf_path, output_path, output_format, pages)
            
            else:
                raise ValueError(f"Unsupported output format: {output_format}")
                
        except Exception as e:
            logger.error(f"PDF conversion error: {e}")
            raise
    
    async def _convert_to_pdf(self, input_path: Path, output_format: str,
                            quality: str) -> Path:
        """Convert other formats to PDF"""
        
        if output_format != "pdf":
            raise ValueError("Can only convert to PDF from other formats")
        
        output_dir = Path("processed")
        output_path = output_dir / f"{input_path.stem}_converted.pdf"
        
        try:
            if input_path.suffix.lower() in [".docx", ".doc"]:
                return await self._word_to_pdf(input_path, output_path)
            
            elif input_path.suffix.lower() in [".xlsx", ".xls"]:
                return await self._excel_to_pdf(input_path, output_path)
            
            elif input_path.suffix.lower() in [".pptx", ".ppt"]:
                return await self._powerpoint_to_pdf(input_path, output_path)
            
            elif input_path.suffix.lower() in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
                return await self._image_to_pdf(input_path, output_path, quality)
            
            elif input_path.suffix.lower() in [".txt", ".html"]:
                return await self._text_to_pdf(input_path, output_path)
            
            else:
                raise ValueError(f"Unsupported input format: {input_path.suffix}")
                
        except Exception as e:
            logger.error(f"To-PDF conversion error: {e}")
            raise
    
    async def _pdf_to_word(self, pdf_path: Path, output_path: Path, 
                          pages: Optional[List[int]]) -> Path:
        """Convert PDF to Word document"""
        try:
            from pdf2docx import Converter
            
            cv = Converter(str(pdf_path))
            cv.convert(str(output_path), start=0, end=None)
            cv.close()
            
            return output_path
            
        except ImportError:
            # Fallback to PyMuPDF
            import fitz
            from docx import Document
            from docx.shared import Inches
            
            doc = Document()
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                if pages and (page_num + 1) not in pages:
                    continue
                
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    doc.add_paragraph(text)
                doc.add_page_break()
            
            doc.save(output_path)
            pdf_document.close()
            
            return output_path
    
    async def _pdf_to_excel(self, pdf_path: Path, output_path: Path,
                           pages: Optional[List[int]]) -> Path:
        """Convert PDF to Excel"""
        try:
            import fitz
            from openpyxl import Workbook
            
            wb = Workbook()
            ws = wb.active
            ws.title = "PDF Data"
            
            pdf_document = fitz.open(pdf_path)
            row = 1
            
            for page_num in range(len(pdf_document)):
                if pages and (page_num + 1) not in pages:
                    continue
                
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=row, column=1, value=line)
                        row += 1
            
            wb.save(output_path)
            pdf_document.close()
            
            return output_path
            
        except Exception as e:
            logger.error(f"PDF to Excel error: {e}")
            raise
    
    async def _pdf_to_powerpoint(self, pdf_path: Path, output_path: Path,
                                pages: Optional[List[int]]) -> Path:
        """Convert PDF to PowerPoint"""
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                if pages and (page_num + 1) not in pages:
                    continue
                
                slide = prs.slides.add_slide(blank_slide_layout)
                page = pdf_document.load_page(page_num)
                
                # Add page number
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), 
                                                Inches(9), Inches(0.5))
                tf = txBox.text_frame
                tf.text = f"Page {page_num + 1}"
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.5), 
                                                     Inches(8), Inches(5))
                    tf2 = txBox2.text_frame
                    tf2.text = text
            
            prs.save(output_path)
            pdf_document.close()
            
            return output_path
            
        except Exception as e:
            logger.error(f"PDF to PowerPoint error: {e}")
            raise
    
    async def _pdf_to_image(self, pdf_path: Path, output_path: Path,
                           format: str, quality: str, pages: Optional[List[int]]) -> Path:
        """Convert PDF pages to images"""
        try:
            import fitz
            from PIL import Image
            
            pdf_document = fitz.open(pdf_path)
            images = []
            
            dpi = {"high": 300, "medium": 150, "low": 72}.get(quality, 150)
            
            for page_num in range(len(pdf_document)):
                if pages and (page_num + 1) not in pages:
                    continue
                
                page = pdf_document.load_page(page_num)
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                
                img_data = pix.tobytes(format)
                img = Image.frombytes("RGB", [pix.width, pix.height], img_data)
                
                if len(pdf_document) == 1 or pages:
                    # Single image
                    img.save(output_path, format=format.upper(), 
                            quality=95 if quality == "high" else 85)
                else:
                    # Multiple images
                    img_path = Path("processed") / f"{pdf_path.stem}_page_{page_num+1}.{format}"
                    img.save(img_path, format=format.upper(), 
                            quality=95 if quality == "high" else 85)
                    images.append(img_path)
            
            pdf_document.close()
            
            if images:
                # Create a zip file or return first image
                return images[0] if len(images) == 1 else self._create_image_zip(images)
            
            return output_path
            
        except Exception as e:
            logger.error(f"PDF to image error: {e}")
            raise
    
    async def _pdf_to_text(self, pdf_path: Path, output_path: Path,
                          format: str, pages: Optional[List[int]]) -> Path:
        """Extract text from PDF"""
        try:
            import fitz
            
            pdf_document = fitz.open(pdf_path)
            text_content = []
            
            for page_num in range(len(pdf_document)):
                if pages and (page_num + 1) not in pages:
                    continue
                
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                text_content.append(f"--- Page {page_num + 1} ---\n{text}\n")
            
            pdf_document.close()
            
            if format == "html":
                html_content = "<html><body>"
                for page_text in text_content:
                    html_content += f"<div class='page'>{page_text.replace(chr(10), '<br>')}</div>"
                html_content += "</body></html>"
                
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
            else:
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(text_content))
            
            return output_path
            
        except Exception as e:
            logger.error(f"PDF to text error: {e}")
            raise
    
    async def _word_to_pdf(self, word_path: Path, output_path: Path) -> Path:
        """Convert Word document to PDF"""
        try:
            # Try using python-docx and reportlab
            from docx import Document
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            doc = Document(word_path)
            c = canvas.Canvas(str(output_path), pagesize=letter)
            width, height = letter
            y = height - inch
            
            for paragraph in doc.paragraphs:
                if y < inch:  # New page
                    c.showPage()
                    y = height - inch
                
                text = paragraph.text.strip()
                if text:
                    c.drawString(inch, y, text)
                    y -= 14  # Line height
            
            c.save()
            return output_path
            
        except Exception as e:
            logger.error(f"Word to PDF error: {e}")
            raise
    
    async def _excel_to_pdf(self, excel_path: Path, output_path: Path) -> Path:
        """Convert Excel to PDF"""
        try:
            from openpyxl import load_workbook
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            wb = load_workbook(excel_path, data_only=True)
            ws = wb.active
            
            c = canvas.Canvas(str(output_path), pagesize=letter)
            width, height = letter
            
            # Simple table rendering
            row_height = 20
            col_width = 80
            start_x = 50
            start_y = height - 50
            
            for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                for c_idx, cell in enumerate(row):
                    x = start_x + (c_idx * col_width)
                    y = start_y - (r_idx * row_height)
                    
                    if y < 50:  # New page
                        c.showPage()
                        start_y = height - 50
                        y = start_y - ((r_idx % 30) * row_height)
                    
                    cell_text = str(cell) if cell is not None else ""
                    c.drawString(x, y, cell_text[:20])  # Limit text length
            
            c.save()
            return output_path
            
        except Exception as e:
            logger.error(f"Excel to PDF error: {e}")
            raise
    
    async def _powerpoint_to_pdf(self, ppt_path: Path, output_path: Path) -> Path:
        """Convert PowerPoint to PDF"""
        try:
            # Simple conversion - extract text to PDF
            from pptx import Presentation
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            prs = Presentation(ppt_path)
            c = canvas.Canvas(str(output_path), pagesize=letter)
            width, height = letter
            
            y = height - 50
            
            for slide_num, slide in enumerate(prs.slides):
                if y < 50:
                    c.showPage()
                    y = height - 50
                
                c.drawString(50, y, f"Slide {slide_num + 1}")
                y -= 20
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            c.drawString(70, y, text[:100])  # Limit text
                            y -= 15
                
                y -= 20
            
            c.save()
            return output_path
            
        except Exception as e:
            logger.error(f"PowerPoint to PDF error: {e}")
            raise
    
    async def _image_to_pdf(self, image_path: Path, output_path: Path, 
                           quality: str) -> Path:
        """Convert image to PDF"""
        try:
            from PIL import Image
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.utils import ImageReader
            
            img = Image.open(image_path)
            c = canvas.Canvas(str(output_path), pagesize=letter)
            width, height = letter
            
            # Calculate scaling
            img_width, img_height = img.size
            scale = min(width/img_width, height/img_height) * 0.8
            
            x = (width - img_width * scale) / 2
            y = (height - img_height * scale) / 2
            
            c.drawImage(ImageReader(img), x, y, 
                       width=img_width * scale, 
                       height=img_height * scale)
            c.save()
            
            return output_path
            
        except Exception as e:
            logger.error(f"Image to PDF error: {e}")
            raise
    
    async def _text_to_pdf(self, text_path: Path, output_path: Path) -> Path:
        """Convert text file to PDF"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            with open(text_path, "r", encoding="utf-8") as f:
                lines = f.readlines()
            
            c = canvas.Canvas(str(output_path), pagesize=letter)
            width, height = letter
            
            y = height - inch
            line_height = 14
            
            for line in lines:
                if y < inch:  # New page
                    c.showPage()
                    y = height - inch
                
                c.drawString(inch, y, line.rstrip())
                y -= line_height
            
            c.save()
            return output_path
            
        except Exception as e:
            logger.error(f"Text to PDF error: {e}")
            raise
    
    def _create_image_zip(self, image_paths: List[Path]) -> Path:
        """Create zip file of multiple images"""
        import zipfile
        
        zip_path = Path("processed") / f"images_{len(image_paths)}.zip"
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for img_path in image_paths:
                zipf.write(img_path, img_path.name)
        
        return zip_path
